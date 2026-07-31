"""课件文本抽取 + 知识点目标自动提取。

支持 pptx / docx / pdf / txt。pdfplumber 对扫描件提取为空时返回空字符串，
调用方据此提示老师手填（避免编造）。
"""
import os
import re

# 目标/知识点 常见动词（命中即视为知识点行）
_GOAL_VERBS = ["掌握", "了解", "认识", "学习", "理解", "能够", "学会", "知道",
               "熟悉", "培养", "体会", "感受", "重点", "难点"]
_BULLET = re.compile(r"^[\s\-\•\·\*\d]+[.、)）]?\s*")


def extract_text(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == ".txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        if ext == ".pptx":
            return _from_pptx(file_path)
        if ext in (".docx", ".doc"):
            return _from_docx(file_path)
        if ext == ".pdf":
            return _from_pdf(file_path)
    except Exception as e:  # noqa: BLE001
        return f"[解析失败：{e}]"
    return ""


def _from_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    out = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    out.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    out.append(" / ".join(c.text for c in row.cells))
    return "\n".join(out)


def _from_docx(path):
    from docx import Document
    doc = Document(path)
    out = [p.text for p in doc.paragraphs if p.text.strip()]
    for tbl in doc.tables:
        for row in tbl.rows:
            out.append(" / ".join(c.text for c in row.cells))
    return "\n".join(out)


def _from_pdf(path):
    import pdfplumber
    out = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            t = page.extract_text() or ""
            if t.strip():
                out.append(t.strip())
    return "\n".join(out)


def extract_objectives(text: str, max_items: int = 8) -> list:
    """从课件文本里挑出可能的知识点目标行（短句 + 含目标动词 或 条目符号）。"""
    if not text:
        return []
    lines = [ln.strip() for ln in re.split(r"[\r\n]+", text) if ln.strip()]
    candidates = []
    for ln in lines:
        if len(ln) > 40:
            continue
        is_bullet = bool(_BULLET.match(ln))
        has_verb = any(v in ln for v in _GOAL_VERBS)
        if is_bullet or has_verb:
            clean = _BULLET.sub("", ln).strip(" 。.、")
            if clean and clean not in candidates:
                candidates.append(clean)
        if len(candidates) >= max_items:
            break
    return candidates
